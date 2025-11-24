#!/usr/bin/env python3
"""
Script to convert Markdown presentation to PowerPoint PPTX
"""

import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def parse_markdown_slides(markdown_content):
    """Parse markdown content into slides"""
    slides = []
    lines = markdown_content.split('\n')

    current_slide = None
    current_content = []

    for line in lines:
        line = line.strip()

        # Check for slide separator (---)
        if line.startswith('---'):
            if current_slide:
                slides.append({
                    'title': current_slide,
                    'content': current_content
                })
            current_slide = None
            current_content = []
            continue

        # Check for slide title (# )
        if line.startswith('# '):
            if current_slide:
                slides.append({
                    'title': current_slide,
                    'content': current_content
                })
            current_slide = line[2:].strip()
            current_content = []
            continue

        # Check for sub-title (## )
        if line.startswith('## '):
            current_content.append(f"## {line[3:].strip()}")
            continue

        # Regular content
        if line and not line.startswith('```') and not line.startswith('title:') and not line.startswith('author:') and not line.startswith('date:'):
            current_content.append(line)

    # Add the last slide
    if current_slide:
        slides.append({
            'title': current_slide,
            'content': current_content
        })

    return slides

def create_slide(prs, slide_layout, title, content_lines):
    """Create a slide with title and content"""
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_placeholder = slide.shapes.title
    title_placeholder.text = title

    # Content
    if content_lines:
        content_text = '\n'.join(content_lines)
        # Remove markdown formatting for display
        content_text = re.sub(r'\*\*(.*?)\*\*', r'\1', content_text)  # Bold
        content_text = re.sub(r'\*(.*?)\*', r'\1', content_text)      # Italic
        content_text = re.sub(r'`(.*?)`', r'\1', content_text)        # Code
        content_text = re.sub(r'^- ', '', content_text, flags=re.MULTILINE)  # Bullet points

        # Try to add content to the first text box
        try:
            content_placeholder = slide.placeholders[1]
            content_placeholder.text = content_text
        except:
            # If no content placeholder, add a text box
            left = Inches(1)
            top = Inches(2)
            width = Inches(8)
            height = Inches(5)
            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.text = content_text

def main():
    # Read markdown file
    try:
        with open('presentation.md', 'r', encoding='utf-8') as f:
            markdown_content = f.read()
    except FileNotFoundError:
        print("Error: presentation.md not found!")
        return

    # Parse slides
    slides = parse_markdown_slides(markdown_content)

    # Create presentation
    prs = Presentation()

    # Slide layouts
    title_slide_layout = prs.slide_layouts[0]  # Title slide
    content_slide_layout = prs.slide_layouts[1]  # Title and content

    # Create slides
    for i, slide_data in enumerate(slides):
        if i == 0:
            # First slide as title slide
            slide = prs.slides.add_slide(title_slide_layout)
            slide.shapes.title.text = slide_data['title']
            if slide_data['content']:
                try:
                    slide.placeholders[1].text = '\n'.join(slide_data['content'])
                except:
                    pass
        else:
            # Content slides
            create_slide(prs, content_slide_layout, slide_data['title'], slide_data['content'])

    # Save presentation
    prs.save('SmulzLab_Presentation.pptx')
    print("Presentation created successfully: SmulzLab_Presentation.pptx")

if __name__ == "__main__":
    main()
